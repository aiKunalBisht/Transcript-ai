"""
exporters/pptx_builder.py — TranscriptAI Enterprise Edition v2
Design: Midnight Executive palette — dark cover/last, light content (sandwich)
Motif: numbered circle icons + severity color bars on every content slide

v2 CHANGE: rendering is now dispatched by slide.slide_type (set by
slide_architect.py v2) instead of by slide position. Three new layouts:
  - said_vs_meant : two-column "what was said" / "what it meant" table,
    color-coded by severity, built straight from soft_rejection_detector
    output via slide_architect — this is the slide the whole deck exists
    for, so it gets its own dedicated treatment rather than the generic
    icon-row layout.
  - decisions     : commitments grouped with a side badge (Japan side /
    Our team / Joint / Unclear) plus owner/deadline, so nothing about who
    owns what gets lost in a flat bullet list.
  - risk_watch    : a small risk register — severity bar + flag + detail,
    sorted highest-severity first.
The cover now also carries a status badge (ON TRACK / NEEDS FOLLOW-UP /
ESCALATE) computed deterministically upstream, so the headline can never
contradict the risk slide underneath it.

_two_col_slide / _icon_row_slide are kept as fallbacks for any slide whose
slide_type isn't recognized (e.g. an older plan dict without slide_type).

NO accent stripes for decoration. Color bars that appear (severity, side
badges) carry meaning, not decoration. Safe fonts: Cambria titles, Calibri
body. Margins=0 on all textboxes.
"""
import io
import logging
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

log = logging.getLogger(__name__)

# ── Midnight Executive palette ────────────────────────────────────────────────
DARK_BG     = RGBColor(0x1E, 0x27, 0x61)   # deep navy
DARK_CARD   = RGBColor(0x28, 0x33, 0x78)   # slightly lighter navy for cards
ICE_BLUE    = RGBColor(0xCA, 0xDC, 0xFC)   # ice blue — accent on dark slides
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE   = RGBColor(0xF8, 0xF9, 0xFF)   # content slide bg
CARD_BG     = RGBColor(0xEE, 0xF2, 0xFF)   # icon circle / card bg
NAVY_TEXT   = RGBColor(0x1E, 0x27, 0x61)   # body text on light slides
MID_TEXT    = RGBColor(0x3D, 0x4E, 0x8A)   # secondary text
SOFT_TEXT   = RGBColor(0x7A, 0x8A, 0xBB)   # captions

# Severity / status colors — these carry meaning (risk level), not decoration.
SEV_HIGH    = RGBColor(0xC0, 0x3B, 0x3B)   # red
SEV_MEDIUM  = RGBColor(0xC9, 0x8A, 0x2E)   # amber
SEV_LOW     = RGBColor(0x5B, 0x7A, 0xAE)   # steel blue
STATUS_OK   = RGBColor(0x2E, 0x8B, 0x57)   # green — ON_TRACK

SIDE_COLORS = {
    "Japan side": RGBColor(0xA8, 0x3A, 0x3A),
    "Our team":   RGBColor(0x1E, 0x6B, 0x4A),
    "Joint":      RGBColor(0x3D, 0x4E, 0x8A),
    "Unclear":    SOFT_TEXT,
}

STATUS_BADGE = {
    "ON_TRACK": ("ON TRACK",        STATUS_OK),
    "WATCH":    ("NEEDS FOLLOW-UP", SEV_MEDIUM),
    "ESCALATE": ("ESCALATE",        SEV_HIGH),
}

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

_PML_NS = "{http://schemas.openxmlformats.org/presentationml/2006/main}"


# ── Safe attribute helpers ────────────────────────────────────────────────────

def _get(obj, attr, default=None):
    """getattr that also works on dicts."""
    if isinstance(obj, dict):
        return obj.get(attr, default)
    return getattr(obj, attr, default)


# ── Plan normalisation ────────────────────────────────────────────────────────

class _SaidVsMeantData:
    __slots__ = ("speaker", "said", "reading", "meant", "severity")

    def __init__(self, src):
        self.speaker  = str(_get(src, "speaker", "Unknown") or "Unknown")
        self.said     = str(_get(src, "said", "") or "")
        self.reading  = str(_get(src, "reading", "") or "")
        self.meant    = str(_get(src, "meant", "") or "")
        self.severity = str(_get(src, "severity", "MEDIUM") or "MEDIUM").upper()


class _CommitmentData:
    __slots__ = ("side", "text", "owner", "deadline")

    def __init__(self, src):
        self.side     = str(_get(src, "side", "Unclear") or "Unclear")
        self.text     = str(_get(src, "text", "") or "")
        self.owner    = str(_get(src, "owner", "") or "")
        self.deadline = str(_get(src, "deadline", "") or "")


class _WatchItemData:
    __slots__ = ("flag", "detail", "severity")

    def __init__(self, src):
        self.flag     = str(_get(src, "flag", "") or "")
        self.detail   = str(_get(src, "detail", "") or "")
        self.severity = str(_get(src, "severity", "MEDIUM") or "MEDIUM").upper()


class _SlideData:
    """
    Lightweight wrapper so slide_data.attr always works regardless of whether
    the upstream source was a Pydantic model, a plain dict, or a dataclass.
    """
    __slots__ = (
        "slide_number", "slide_type", "title", "bullets", "speaker_notes",
        "estimated_duration_seconds", "said_vs_meant", "commitments", "watch_items",
    )

    def __init__(self, src):
        self.slide_number   = int(_get(src, "slide_number", 1))
        self.slide_type      = str(_get(src, "slide_type") or "content")
        self.title            = str(_get(src, "title", "Slide"))
        raw_bullets           = _get(src, "bullets") or []
        self.bullets          = [str(b) for b in raw_bullets if b]
        self.speaker_notes    = str(_get(src, "speaker_notes") or "")
        self.estimated_duration_seconds = int(
            _get(src, "estimated_duration_seconds") or 90
        )
        self.said_vs_meant = [_SaidVsMeantData(x) for x in (_get(src, "said_vs_meant") or [])]
        self.commitments   = [_CommitmentData(x) for x in (_get(src, "commitments") or [])]
        self.watch_items    = [_WatchItemData(x) for x in (_get(src, "watch_items") or [])]


class _Plan:
    """
    Normalised presentation plan — accepts a PresentationPlan Pydantic model,
    a plain dict (JSON round-trip), or any object with the right attributes.
    """
    def __init__(self, src):
        self.meeting_title     = str(_get(src, "meeting_title") or "Meeting")
        self.status_flag        = str(_get(src, "status_flag") or "ON_TRACK").upper()
        self.meeting_context    = str(_get(src, "meeting_context") or "")
        self.executive_summary  = str(_get(src, "executive_summary") or "")
        self.language            = str(_get(src, "language") or "en")
        raw_slides               = _get(src, "slides") or []
        self.slides              = [_SlideData(s) for s in raw_slides]
        self.total_slides        = len(self.slides)


# ── Low-level drawing helpers ─────────────────────────────────────────────────

def _bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _tb(slide, l, t, w, h):
    shape = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = True
    return tf


def _run(tf, text, size, color, font="Calibri", bold=False, italic=False,
         align=PP_ALIGN.LEFT):
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.name = font
    r.font.bold = bold
    r.font.italic = italic
    return r


def _strip_theme_style(shape):
    """Remove <p:style> so LibreOffice doesn't apply phantom drop-shadows."""
    sp = shape._element
    style = sp.find(_PML_NS + "style")
    if style is not None:
        sp.remove(style)


def _rect(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.shadow.inherit = False
    _strip_theme_style(s)
    return s


def _circle(slide, l, t, d, color):
    s = slide.shapes.add_shape(9, Inches(l), Inches(t), Inches(d), Inches(d))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.shadow.inherit = False
    _strip_theme_style(s)
    return s


def _num_in_circle(slide, number, l, t, d=0.44, bg=None, fg=None):
    if bg is None:
        bg = DARK_BG
    if fg is None:
        fg = WHITE
    _circle(slide, l, t, d, bg)
    tf = _tb(slide, l, t + 0.02, d, d - 0.04)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = tf.paragraphs[0].add_run()
    r.text = str(number)
    r.font.size = Pt(13)
    r.font.bold = True
    r.font.color.rgb = fg
    r.font.name = "Calibri"


def _slide_num(slide, n, total, fg=None):
    if fg is None:
        fg = SOFT_TEXT
    tf = _tb(slide, 12.3, 7.1, 0.9, 0.3)
    tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
    r = tf.paragraphs[0].add_run()
    r.text = f"{n} / {total}"
    r.font.size = Pt(9)
    r.font.color.rgb = fg
    r.font.name = "Calibri"


def _section_header(slide, title):
    """Dark title bar used by content-heavy slides (said_vs_meant, risk_watch)."""
    _rect(slide, 0.4, 0.3, 12.53, 0.95, DARK_BG)
    tf_t = _tb(slide, 0.7, 0.52, 11.9, 0.6)
    _run(tf_t, title, 24, WHITE, font="Cambria", bold=True)


def _light_header(slide, title):
    """Light underlined title used by data-table slides (decisions, bottom_line)."""
    tf_t = _tb(slide, 0.4, 0.22, 12.5, 0.9)
    _run(tf_t, title, 30, NAVY_TEXT, font="Cambria", bold=True)
    _rect(slide, 0.4, 1.18, 12.5, 0.04, RGBColor(0xCC, 0xD4, 0xF0))


def _row_layout(start_y, band_h, n, ideal_row_h=1.6, min_row_h=0.9):
    row_h = max(min_row_h, min(ideal_row_h, band_h / max(n, 1)))
    content_h = row_h * n
    if content_h < band_h:
        start_y = start_y + (band_h - content_h) / 2
    else:
        row_h = band_h / n
    return start_y, row_h


def _sev_color(severity):
    return {"HIGH": SEV_HIGH, "MEDIUM": SEV_MEDIUM, "LOW": SEV_LOW}.get(
        (severity or "MEDIUM").upper(), SEV_MEDIUM
    )


# ══════════════════════════════════════════════════════════════════════════════
# COVER — dark, with status badge + meeting context
# ══════════════════════════════════════════════════════════════════════════════
def _cover_slide(prs, plan):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, DARK_BG)
    _rect(slide, 7.8, 0, 5.53, 7.5, DARK_CARD)

    tf = _tb(slide, 0.5, 0.38, 5.4, 0.3)
    _run(tf, "TRANSCRIPTAI  ·  MEETING INTELLIGENCE", 8, ICE_BLUE, bold=True)

    # Status badge — top right of the light column, computed deterministically
    # upstream so it can never disagree with the Risk & Watch slide.
    label, color = STATUS_BADGE.get(plan.status_flag, STATUS_BADGE["ON_TRACK"])
    badge_w = max(1.5, 0.09 * len(label) + 0.5)
    _rect(slide, 7.8 - badge_w - 0.25, 0.36, badge_w, 0.36, color)
    tf_s = _tb(slide, 7.8 - badge_w - 0.25, 0.43, badge_w, 0.25)
    _run(tf_s, label, 9, WHITE, bold=True, align=PP_ALIGN.CENTER)

    # Main title
    tf_t = _tb(slide, 0.5, 0.85, 7.0, 1.8)
    r = tf_t.paragraphs[0].add_run()
    r.text = plan.meeting_title
    r.font.size = Pt(36)
    r.font.bold = True
    r.font.color.rgb = WHITE
    r.font.name = "Cambria"

    # Meeting context — who was there, why it mattered
    y = 2.75
    if plan.meeting_context:
        tf_ctx = _tb(slide, 0.5, y, 7.0, 0.8)
        _run(tf_ctx, plan.meeting_context, 13, ICE_BLUE, italic=True)
        y += 0.9

    # Executive summary
    if plan.executive_summary:
        tf_e = _tb(slide, 0.5, y, 7.0, 1.3)
        _run(tf_e, plan.executive_summary[:220], 14, WHITE)
        y += 1.0

    # Stats row
    dur_sec = sum(s.estimated_duration_seconds for s in plan.slides)
    dur_min = max(1, dur_sec // 60)
    stats = [
        (str(plan.total_slides), "SLIDES"),
        (f"{dur_min}m", "DURATION"),
        (plan.language.upper()[:5], "LANGUAGE"),
    ]
    stat_y = 6.2
    for i, (val, lbl) in enumerate(stats):
        x = 0.5 + i * 2.2
        tf_v = _tb(slide, x, stat_y, 2.0, 0.55)
        _run(tf_v, val, 24, WHITE, font="Cambria", bold=True)
        tf_l = _tb(slide, x, stat_y + 0.55, 2.0, 0.3)
        _run(tf_l, lbl, 8, SOFT_TEXT, bold=True)

    # Right panel decorative
    tf_w = _tb(slide, 7.95, 1.2, 5.2, 5.0)
    r_w = tf_w.paragraphs[0].add_run()
    r_w.text = plan.meeting_title
    r_w.font.size = Pt(28)
    r_w.font.bold = True
    r_w.font.color.rgb = RGBColor(0x38, 0x45, 0x90)
    r_w.font.name = "Cambria"

    tf_d = _tb(slide, 7.95, 6.5, 5.0, 0.4)
    _run(tf_d, datetime.now().strftime("%B %d, %Y"), 11, SOFT_TEXT)

    _slide_num(slide, 1, plan.total_slides, fg=SOFT_TEXT)


# ══════════════════════════════════════════════════════════════════════════════
# BOTTOM LINE — light, numbered rows
# ══════════════════════════════════════════════════════════════════════════════
def _bottom_line_slide(prs, slide_data, plan):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, OFF_WHITE)
    _light_header(slide, slide_data.title)

    bullets = (slide_data.bullets or ["Details unavailable."])[:4]
    n = len(bullets)
    font_sz = {1: 22, 2: 19, 3: 17, 4: 16}.get(n, 15)
    start_y, row_h = _row_layout(1.5, 5.1, n, ideal_row_h=1.6)

    for idx, b in enumerate(bullets):
        y = start_y + idx * row_h
        _num_in_circle(slide, idx + 1, 0.5, y + 0.05, 0.42, DARK_BG, WHITE)
        tf_b = _tb(slide, 1.1, y, 11.2, row_h - 0.1)
        _run(tf_b, b, font_sz, NAVY_TEXT)

    if slide_data.speaker_notes:
        tf_note = _tb(slide, 0.5, 6.78, 12.3, 0.55)
        _run(tf_note, slide_data.speaker_notes[:170], 10, MID_TEXT, italic=True)

    _slide_num(slide, slide_data.slide_number, plan.total_slides)


# ══════════════════════════════════════════════════════════════════════════════
# SAID VS MEANT — two-column table, severity-coded
# ══════════════════════════════════════════════════════════════════════════════
def _said_vs_meant_slide(prs, slide_data, plan):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, WHITE)
    _section_header(slide, slide_data.title)

    tf_h1 = _tb(slide, 0.5, 1.45, 5.6, 0.3)
    _run(tf_h1, "WHAT WAS SAID", 9, SOFT_TEXT, bold=True)
    tf_h2 = _tb(slide, 6.5, 1.45, 6.3, 0.3)
    _run(tf_h2, "WHAT IT MEANT", 9, SOFT_TEXT, bold=True)
    _rect(slide, 0.5, 1.78, 12.3, 0.02, RGBColor(0xCC, 0xD4, 0xF0))

    items = (slide_data.said_vs_meant or [])[:4]
    n = len(items) or 1
    start_y, row_h = _row_layout(1.95, 4.95, n, ideal_row_h=1.55, min_row_h=1.05)

    for idx, item in enumerate(items):
        y = start_y + idx * row_h
        color = _sev_color(item.severity)
        _rect(slide, 0.5, y, 0.07, row_h - 0.18, color)

        tf_spk = _tb(slide, 0.7, y, 5.3, 0.28)
        _run(tf_spk, item.speaker.upper(), 9, MID_TEXT, bold=True)

        tf_said = _tb(slide, 0.7, y + 0.3, 5.3, row_h - 0.45)
        p = tf_said.paragraphs[0]
        r = p.add_run()
        r.text = item.said or "\u2014"
        r.font.size = Pt(14)
        r.font.bold = True
        r.font.color.rgb = NAVY_TEXT
        r.font.name = "Calibri"
        if item.reading:
            p2 = tf_said.add_paragraph()
            r2 = p2.add_run()
            r2.text = f"\u201c{item.reading}\u201d"
            r2.font.size = Pt(10)
            r2.font.italic = True
            r2.font.color.rgb = SOFT_TEXT

        tf_meant = _tb(slide, 6.5, y, 6.3, row_h - 0.18)
        _run(tf_meant, item.meant or "\u2014", 12, NAVY_TEXT)

    _slide_num(slide, slide_data.slide_number, plan.total_slides)


# ══════════════════════════════════════════════════════════════════════════════
# DECISIONS & COMMITMENTS — side-badged rows
# ══════════════════════════════════════════════════════════════════════════════
def _decisions_slide(prs, slide_data, plan):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, OFF_WHITE)
    _light_header(slide, slide_data.title)

    items = (slide_data.commitments or [])[:4]
    n = len(items) or 1
    start_y, row_h = _row_layout(1.4, 5.5, n, ideal_row_h=1.55, min_row_h=1.0)

    for idx, item in enumerate(items):
        y = start_y + idx * row_h
        _rect(slide, 0.4, y, 12.5, row_h - 0.15, CARD_BG)

        color = SIDE_COLORS.get(item.side, SOFT_TEXT)
        badge_w = max(1.3, 0.09 * len(item.side) + 0.4)
        _rect(slide, 0.6, y + 0.14, badge_w, 0.32, color)
        tf_badge = _tb(slide, 0.6, y + 0.2, badge_w, 0.22)
        _run(tf_badge, item.side.upper(), 8, WHITE, bold=True, align=PP_ALIGN.CENTER)

        tf_txt = _tb(slide, 0.6, y + 0.56, 9.6, row_h - 0.72)
        _run(tf_txt, item.text or "\u2014", 14, NAVY_TEXT)

        meta = " \u00b7 ".join(filter(None, [item.owner, item.deadline]))
        if meta:
            tf_meta = _tb(slide, 10.3, y + 0.16, 2.3, 0.3)
            _run(tf_meta, meta, 9, MID_TEXT, align=PP_ALIGN.RIGHT)

    _slide_num(slide, slide_data.slide_number, plan.total_slides)


# ══════════════════════════════════════════════════════════════════════════════
# RISK & WATCH ITEMS — severity-coded register
# ══════════════════════════════════════════════════════════════════════════════
def _risk_watch_slide(prs, slide_data, plan):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, WHITE)
    _section_header(slide, slide_data.title)

    items = (slide_data.watch_items or [])[:4]
    n = len(items) or 1
    start_y, row_h = _row_layout(1.45, 5.6, n, ideal_row_h=1.6, min_row_h=1.0)

    for idx, item in enumerate(items):
        y = start_y + idx * row_h
        color = _sev_color(item.severity)
        if idx % 2 == 0:
            _rect(slide, 0.3, y + 0.04, 12.7, row_h - 0.08, RGBColor(0xF4, 0xF6, 0xFF))
        _rect(slide, 0.45, y + 0.12, 0.1, row_h - 0.34, color)

        tf_flag = _tb(slide, 0.78, y + 0.1, 11.8, 0.4)
        _run(tf_flag, item.flag or "Flagged", 15, NAVY_TEXT, bold=True)
        tf_detail = _tb(slide, 0.78, y + 0.52, 11.6, row_h - 0.65)
        _run(tf_detail, item.detail or "", 12, MID_TEXT)

    _slide_num(slide, slide_data.slide_number, plan.total_slides)


# ══════════════════════════════════════════════════════════════════════════════
# FALLBACKS — generic content layouts for unrecognized/legacy slide_type
# ══════════════════════════════════════════════════════════════════════════════
def _two_col_slide(prs, slide_data, plan):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, OFF_WHITE)
    _light_header(slide, slide_data.title)

    bullets = (slide_data.bullets or ["Details in transcript"])[:4]
    n = len(bullets)
    font_sz = {1: 22, 2: 19, 3: 16, 4: 15}.get(n, 15)

    start_y, row_h = _row_layout(1.4, 5.5, n, ideal_row_h=1.7)
    for b_idx, bullet in enumerate(bullets):
        y = start_y + b_idx * row_h
        _num_in_circle(slide, b_idx + 1, 0.4, y + 0.05, 0.42, DARK_BG, WHITE)
        tf_b = _tb(slide, 0.98, y, 6.8, row_h - 0.1)
        _run(tf_b, bullet, font_sz, NAVY_TEXT)

    _rect(slide, 7.8, 1.35, 5.1, 5.7, CARD_BG)
    tf_big = _tb(slide, 8.0, 1.55, 4.7, 1.6)
    _run(tf_big, f"0{slide_data.slide_number}", 72, DARK_BG,
         font="Cambria", bold=True, align=PP_ALIGN.CENTER)

    tf_sub = _tb(slide, 8.0, 3.1, 4.7, 0.5)
    _run(tf_sub, slide_data.title.upper()[:50], 9, SOFT_TEXT,
         bold=True, align=PP_ALIGN.CENTER)

    notes_preview = slide_data.speaker_notes[:140] if slide_data.speaker_notes else ""
    if notes_preview:
        tf_np = _tb(slide, 8.1, 3.75, 4.6, 2.8)
        _run(tf_np, f'"{notes_preview}"', 11, MID_TEXT, italic=True)

    _slide_num(slide, slide_data.slide_number, plan.total_slides)


def _icon_row_slide(prs, slide_data, plan):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, WHITE)
    _section_header(slide, slide_data.title)

    bullets = (slide_data.bullets or ["Details in transcript"])[:5]
    n = len(bullets)
    font_sz = {1: 20, 2: 18, 3: 16, 4: 15, 5: 14}.get(n, 14)
    ICON_CHARS = ["\u25c6", "\u25b8", "\u25cf", "\u2605", "\u25c9"]

    start_y, row_h = _row_layout(1.45, 5.6, n, ideal_row_h=1.6)
    for b_idx, bullet in enumerate(bullets):
        y = start_y + b_idx * row_h
        if b_idx % 2 == 0:
            _rect(slide, 0.3, y + 0.04, 12.7, row_h - 0.08, RGBColor(0xF4, 0xF6, 0xFF))
        _num_in_circle(slide, b_idx + 1, 0.45, y + (row_h - 0.44) / 2,
                       0.44, DARK_BG, ICE_BLUE)
        tf_ic = _tb(slide, 1.05, y + (row_h - 0.44) / 2, 0.44, 0.44)
        _run(tf_ic, ICON_CHARS[b_idx % len(ICON_CHARS)], 14,
             ICE_BLUE, align=PP_ALIGN.CENTER)
        tf_b = _tb(slide, 1.55, y + (row_h - 0.55) / 2, 11.4, row_h * 0.85)
        _run(tf_b, bullet, font_sz, NAVY_TEXT)

    _slide_num(slide, slide_data.slide_number, plan.total_slides)


# ══════════════════════════════════════════════════════════════════════════════
# CLOSING — dark, next steps as numbered pills
# ══════════════════════════════════════════════════════════════════════════════
def _closing_slide(prs, slide_data, plan):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, DARK_BG)
    _rect(slide, 8.5, 0, 4.83, 7.5, DARK_CARD)

    tf_t = _tb(slide, 0.5, 0.35, 7.7, 1.0)
    _run(tf_t, slide_data.title, 32, WHITE, font="Cambria", bold=True)
    tf_s = _tb(slide, 0.5, 1.35, 7.7, 0.4)
    _run(tf_s, "NEXT STEPS  \u00b7  ACTION REQUIRED", 9, ICE_BLUE, bold=True)

    bullets = (slide_data.bullets or ["Follow up with all stakeholders"])[:5]
    n = len(bullets)
    font_sz = {1: 20, 2: 18, 3: 15, 4: 14, 5: 13}.get(n, 13)

    start_y, row_h = _row_layout(1.9, 4.8, n, ideal_row_h=1.5, min_row_h=0.85)
    for b_idx, bullet in enumerate(bullets):
        y = start_y + b_idx * row_h
        _rect(slide, 0.5, y, 7.7, row_h - 0.12, RGBColor(0x28, 0x33, 0x78))
        _num_in_circle(slide, b_idx + 1, 0.6, y + 0.06, 0.42, ICE_BLUE, DARK_BG)
        tf_b = _tb(slide, 1.18, y + 0.06, 6.8, row_h - 0.18)
        _run(tf_b, bullet, font_sz, WHITE)

    # Right panel brand
    tf_br = _tb(slide, 8.65, 0.6, 4.5, 1.2)
    _run(tf_br, "TranscriptAI", 24, ICE_BLUE, font="Cambria",
         bold=True, align=PP_ALIGN.CENTER)
    tf_bl = _tb(slide, 8.65, 1.75, 4.5, 0.35)
    _run(tf_bl, "Meeting Intelligence Platform", 10, SOFT_TEXT,
         align=PP_ALIGN.CENTER)
    tf_g = _tb(slide, 8.65, 6.88, 4.5, 0.35)
    _run(tf_g, "github.com/aiKunalBisht/Transcript-ai", 8,
         SOFT_TEXT, align=PP_ALIGN.CENTER)

    _slide_num(slide, slide_data.slide_number, plan.total_slides, fg=SOFT_TEXT)


# ══════════════════════════════════════════════════════════════════════════════
# PUBLIC API
# ══════════════════════════════════════════════════════════════════════════════
_RENDERERS = {
    "bottom_line":     _bottom_line_slide,
    "said_vs_meant":   _said_vs_meant_slide,
    "decisions":       _decisions_slide,
    "risk_watch":      _risk_watch_slide,
    "closing":         _closing_slide,
}


def build_pptx(plan_input) -> bytes:
    """
    Builds enterprise PPTX from a PresentationPlan, plain dict, or any object
    with the expected attributes. Normalises input through _Plan so that JSON
    round-trips, Pydantic models, and raw dicts all work identically.

    Slides are dispatched by slide.slide_type. "cover" is handled specially
    (it reads from `plan`, not `slide_data`). Any unrecognized slide_type
    (e.g. a legacy plan without slide_type, or "content") falls back to
    _two_col_slide for the first such slide and _icon_row_slide thereafter,
    matching the original v1 layout behavior.

    Returns raw bytes. Never raises — returns a minimal error slide on failure.
    """
    try:
        plan = _Plan(plan_input)

        if plan.total_slides == 0:
            raise ValueError("PresentationPlan has no slides to render.")

        prs = Presentation()
        prs.slide_width  = SLIDE_W
        prs.slide_height = SLIDE_H

        fallback_used = 0
        for slide_data in plan.slides:
            stype = slide_data.slide_type

            if stype == "cover":
                _cover_slide(prs, plan)
            elif stype in _RENDERERS:
                _RENDERERS[stype](prs, slide_data, plan)
            else:
                if fallback_used == 0:
                    _two_col_slide(prs, slide_data, plan)
                else:
                    _icon_row_slide(prs, slide_data, plan)
                fallback_used += 1

            # Speaker notes
            try:
                notes_slide = prs.slides[-1].notes_slide
                notes_slide.notes_text_frame.text = slide_data.speaker_notes or ""
            except Exception:
                pass

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        return buf.read()

    except Exception as e:
        log.error("pptx_builder failed: %s", e, exc_info=True)
        try:
            prs = Presentation()
            prs.slide_width  = SLIDE_W
            prs.slide_height = SLIDE_H
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = DARK_BG
            tb = slide.shapes.add_textbox(
                Inches(0.5), Inches(2.5), Inches(12), Inches(2)
            )
            p = tb.text_frame.paragraphs[0]
            r = p.add_run()
            r.text = (
                f"TranscriptAI — PPTX generation error.\n"
                f"Details: {str(e)[:300]}\n"
                f"Please try exporting again or use the JSON export."
            )
            r.font.size = Pt(16)
            r.font.color.rgb = WHITE
            r.font.name = "Calibri"
            buf = io.BytesIO()
            prs.save(buf)
            buf.seek(0)
            return buf.read()
        except Exception:
            return b""